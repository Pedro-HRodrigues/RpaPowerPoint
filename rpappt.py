from http.client import PROXY_AUTHENTICATION_REQUIRED
import io
from logging import root
import os
import glob
from re import sub
from PIL import Image
import warnings
from numpy import tile
import pyperclip
import win32com.client
from time import sleep
import pyautogui as pyauto
import smtplib
import email.message
import pandas as pd
from tkinter import Tk
from tkinter.messagebox import Message 
from _tkinter import TclError
from pptx import Presentation
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from datetime import datetime

warnings.filterwarnings("ignore")
mesesPortugues = {1:'Janeiro', 2:'Fevereiro', 3:'Março', 4:'Abril', 5:'Maio', 6:'Junho', 7:'Julho', 8:'Agosto', 9:'Setembro', 10:'Outubro', 11:'Novembro', 12:'Dezembro'}
pastaTemp = os.getenv('TEMP')
pastaAtual = os.getcwd()
pyauto.PAUSE = 1

def enviar_email(corpo, titulo):
    arquivo = open(os.path.join(pastaAtual, 'log-rpaPPT.txt'),'a')
    arquivo.write(datetime.now().strftime('%d/%m/%Y, %H:%M:%S')+'- PROBLEMA ' + titulo + ';ERRO ' + sub('[^A-Za-z0-9 ><]', '', corpo) + "\n")
    arquivo.close()
    try:
        msg = email.message.Message()
        msg['Subject'] = titulo
        msg['From'] = 'contato@evope.com.br'
        msg['To'] = 'vinicius@evope.com.br'
        password = '3v0'
        msg.add_header('Content-Type', 'text/html')
        msg.set_payload(corpo)
        s = smtplib.SMTP('smtp.gmail.com', 587)
        s.ehlo()
        s.starttls()
        s.login(msg['From'], password)
        s.sendmail(msg['From'], [msg['To']], msg.as_string().encode('utf-8'))
        s.quit()
    except Exception as ex:
        arquivo = open(os.path.join(pastaAtual, 'log-rpaPPT.txt'),'a')
        arquivo.write(datetime.now().strftime('%d/%m/%Y, %H:%M:%S')+' - ERRO AO ENVIAR E-MAIL REFERENTE AO PROBLEMA ' + titulo + ', ERRO ' + corpo + "\n")
        arquivo.close()

def trocarTextoMantendoFormatacaoNmGrupo(paragraph, novoTexto):
    p = paragraph._p
    for idx, run in enumerate(paragraph.runs):
        if idx == 0: 
            continue
        p.remove(run._r)
    if(len(novoTexto)>=20): 
        paragraph.runs[0].font.size = int(paragraph.runs[0].font.size*(20/len(novoTexto)))
    paragraph.runs[0].text = novoTexto

def trocarTextoMantendoFormatacaoGerencial(paragraph, novoTexto):
    p = paragraph._p
    for idx, run in enumerate(paragraph.runs):
        if idx == 0: 
            continue
        p.remove(run._r)
    if(len(novoTexto)>=16): 
        paragraph.runs[0].font.size = int(paragraph.runs[0].font.size*(24/len(novoTexto)))
    paragraph.runs[0].text = novoTexto
    
def trocarTextoMantendoFormatacao(paragraph, novoTexto):
    p = paragraph._p
    for idx, run in enumerate(paragraph.runs):
        if idx == 0: 
            continue
        p.remove(run._r)
    paragraph.runs[0].text = novoTexto

def paginaCarregada(campo):
    sleep(1)
    paginaCarregando = True
    while paginaCarregando:
        try:
            driver.find_element(By.XPATH, campo).click()
            paginaCarregando = False
        except Exception as ex:
            pyauto.scroll(200)
            sleep(1)
            pyauto.scroll(200)

#importar o arquivo de configuração em modo utf-8
try:
    f = io.open(os.path.join(pastaAtual, 'clientesCSV.csv'), mode='r', encoding='utf-8')
    data = f.read()
    rows = data.split('\n')
    clientes = []
    for row in rows: 
        clientes.append(row.split(";"))
    
    dfClientes = pd.DataFrame(clientes, columns = ['clientes','grupos','tipo'])
    dfClientes = dfClientes.drop(dfClientes.index[0])
    numGrupos = dfClientes.groupby('clientes')
    nomesClientes = dfClientes['clientes'].unique()
    
except Exception as ex:
    corpo = """<h1>Erro Evope - Leitura CSV</h1><p>Favor verificar o arquivo *.CSV com os clientes que o RPA deve atuar.</p><p>Caso o erro persista entrar em contato com a Evope.</p><p>Muito obrigado!<br>Equipe Evope</p>"""
    enviar_email(corpo, 'RPA - Erro Evope - Leitura CSV')
    exit()

#solicitar informações de Início e fim
print('Gerador de PPTX - Evope')
print('Favor informar:')
#solicitar data de início e fim no terminal
dtInicio = input('Data de início, ex.: 01/01/2000  -> ')
dtFim = input('Data de término, ex.: 01/01/2000  -> ')
#dtInicio = '01/07/2022'
#dtFim = '31/07/2022'
print('Obrigado, só aguardar!')

# carregar usuário e senha do portal evope
try:
    f = io.open(os.path.join(pastaAtual, 'config'), mode='r', encoding='utf-8')
    data = f.read()
    username = data.split(';')[0]
    password = data.split(';')[1]
except Exception as ex:
    corpo = """<h1>Erro Evope - Carregar usuário e senha</h1><p>Favor verificar o arquivo config.</p><p>Caso o erro persista entrar em contato com a Evope.</p><p>Muito obrigado!<br>Equipe Evope</p>"""
    enviar_email(corpo, 'RPA - Erro Evope - Usuário e Senha')
    exit()

try:
    #ocultar mensagens do webdriver
    options = webdriver.ChromeOptions()    
    options.add_argument("--incognito")
    #iniciar driver para o Chrome
    driver = webdriver.Chrome(chrome_options=options)
    #fazer login
    driver.get("https://portal.evope.com.br/account/login")
    sleep(2)
    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/section/div/div[1]/input').send_keys(username)
    pyauto.press('enter')
    sleep(1)
    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/section/div/div[1]/div[1]/input').send_keys(password)
    pyauto.press('enter')
    #sleep(0.5)
    driver.maximize_window()
    telaX, telaY = pyauto.size()
    
    paginaCarregada('/html/body/form/div[4]/div/div[2]/div/div/select')
    pyauto.moveTo(300, 300)
except Exception as ex:
    corpo = """<h1>Erro Evope - Erro ao acessar o portal</h1><p>Favor verificar conexão com a internet e executar o RPA novamente.</p><p>Caso o erro persista entrar em contato com a Evope.</p><p>Muito obrigado!<br>Equipe Evope</p>"""
    enviar_email(corpo, 'RPA - Erro Evope - Erro ao acessar o portal')
    exit()

#apagar todas as imagens da pasta para evitar erros
py_files = glob.glob(pastaTemp + '\\*.png')
for py_file in py_files:
    try:
        os.remove(py_file)
    except OSError as e:
        pass
py_files = glob.glob('C:\\apresentacoes' + '\\*.pptx')
for py_file in py_files:
    try:
        # send2trash(py_file)
        os.remove(py_file)
    except OSError as e:
        pass

#carregar as informações por grupos de clientes
pptIntancia = []
for y in range(len(nomesClientes)):
    #pptIntancia.append('ppt_instance' + str(y))
    pptIntancia.append(win32com.client.Dispatch("Powerpoint.Application"))
    
for y in range(len(nomesClientes)):
    #for por clientes para começar a carregar os dados para montar a apresentação
    element = driver.find_element(By.XPATH, '/html/body/form/div[4]/div/div[2]/div/div/select')
    all_li = element.find_elements(By.TAG_NAME, "option")
    for li in all_li:
        if(li.text.upper().count(nomesClientes[y].upper())==1):
            li.click()
            sleep(1.5)
            break
    # #mensagem de aviso de criação do powerpoint
    prsTeste = pptIntancia[y].Presentations.open("C:/apresentacoes/modelo/modeloGerencial.pptx", True, False, False)
    TIME_TO_WAIT = 3000 # in milliseconds 
    root = Tk() 
    root.withdraw()
    try:
        root.after(TIME_TO_WAIT, root.destroy) 
        Message(title="Por favor, Aguarde!", message="A montagem do *.PPTX pode demorar um tempo.", master=root).show()
    except TclError:
        pass
    # arrumar a estrutura da apresentação
    sleep(1)
    for x in range(1, len(numGrupos.groups[nomesClientes[y]])):
        for z in range(2,8):
            prsTeste.Slides(z).Copy()
            sleep(0.5)
            prsTeste.Slides.Paste(Index=z+6)
    prsTeste.SaveAs('C:/apresentacoes/' + nomesClientes[y].replace(' ','')+'-'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2]) + '.pptx')
    prsTeste.Close()
    sleep(2)
    #carregar a apresentação do cliente
    prs = Presentation('C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'-'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.pptx')
    qtdeColaboradoresGrupo = 0
    #em uma guia o mapa de processos e na nova guia o relatório gerencial 
    if(len(driver.window_handles)<=1):
        pyauto.hotkey('ctrl','t')
    driver.switch_to.window(driver.window_handles[1])
    #apenas para teste, depois remover
    sleep(3)
    driver.get('https://portal.evope.com.br/relatoriogerencial')
    paginaCarregada('/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]')
    for a in range(10):
        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(Keys.BACKSPACE)    
        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(Keys.DELETE)
    for a in range(10):
        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(Keys.BACKSPACE)    
        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(Keys.DELETE)            
    sleep(1)
    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(dtInicio)
    sleep(1)
    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(dtFim)
    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/a[1]').click()
    paginaCarregada('/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]')
    qtdeColaboradoresGrupo =  int(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[1]/div[1]/table/tbody/tr[2]/td[3]/label').text)
    driver.switch_to.window(driver.window_handles[0])
    sleep(1)
    #carregar as informações do portal
    for x in range(len(numGrupos.groups[nomesClientes[y]])):
        pyauto.moveTo(500, 500)
        try:
            # mapa de processo
            driver.get("https://portal.evope.com.br/ProcessMap")
            paginaCarregada('/html/body/form/div[5]/div/div[1]/div[1]/div[5]/input')
            paginaCarregando = True
            while paginaCarregando:
                try:
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[1]/input[1]').click()
                    paginaCarregando = False
                except Exception as ex:
                    pyauto.PAUSE = 0.1
                    pyauto.scroll(50)
                    pyauto.press('pageup')
                    pyauto.scroll(50)
                    pyauto.press('pageup')
                    pyauto.scroll(50)
                    pyauto.press('pageup')
                    pyauto.PAUSE = 1
            pyauto.hotkey('ctrl','a')
            pyauto.write(dtInicio, interval=0.02)
            pyauto.press('tab')
            pyauto.write(dtFim, interval=0.02)
            #selecionar o grupo, grupos ou colaborador
            if(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '2'):
                try:
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[4]/span/div/button').click()
                    sleep(0.4)
                    element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[4]/span/div/ul')
                except Exception as ex:
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/button').click()
                    sleep(0.4)
                    element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/ul')
                sleep(0.5)
                all_li = element.find_elements(By.TAG_NAME, "li")
                grupoNaoSelecionado = True
                for li in all_li:
                    if(li.text.upper().count(clientes[numGrupos.groups[nomesClientes[y]][x]][1].upper())==1):
                            li.click()
                            grupoNaoSelecionado = False
                            break
            elif(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '3'):
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/button').click()
                sleep(1)
                element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/ul')
                all_li = element.find_elements(By.TAG_NAME, "li")
                grupoNaoSelecionado = True
                agrupados = clientes[numGrupos.groups[nomesClientes[y]][x]][1].split(',')
                for li in all_li:                    
                    for numAgrupados in range(len(agrupados)):
                        if(li.text.upper().count('MASTER > ' + agrupados[numAgrupados].upper())==1):
                            li.click()
                            grupoNaoSelecionado = False
                            break
            else:
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/button').click()
                sleep(1)
                element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[3]/span/div/ul')
                all_li = element.find_elements(By.TAG_NAME, "li")
                grupoNaoSelecionado = True
                if(clientes[numGrupos.groups[nomesClientes[y]][x]][1] == 'tudo'):
                    for li in all_li:
                        if(li.text.upper().count('MASTER')==1):
                            li.click()
                            grupoNaoSelecionado = False
                            break
                elif(clientes[numGrupos.groups[nomesClientes[y]][x]][1] == 'geral'):
                    for li in all_li:
                        if(li.text.upper().count('SELECT ALL')==1):
                            li.click()
                            grupoNaoSelecionado = False
                            break
                else:
                    for li in all_li:
                        if(li.text.upper().count('MASTER > ' + clientes[numGrupos.groups[nomesClientes[y]][x]][1].upper())==1):
                            li.click()         
                            grupoNaoSelecionado = False
                            break
            pyauto.PAUSE = 1
            if(grupoNaoSelecionado):
                corpo = """<h1>Erro Evope - Erro ao selecionar o cliente</h1><p>Cliente {} e Grupo {}<br>Favor verificar o config e executar o RPA navamente.</p><p>Caso o erro persista entrar em contato com a Evope.</p><p>Muito obrigado!<br>Equipe Evope</p>""".format(clientes[x+1][0],clientes[x+1][1])
                enviar_email(corpo, 'RPA - Erro Evope - Erro ao selecionar o cliente')    
            else:
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div[1]/div[5]/input').click()
                sleep(2)
                paginaCarregada('/html/body/form/div[5]/div/div[1]/div[1]/div[5]/input')
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div/button').click()
                paginaCarregada('/html/body/form/div[5]/div/div[1]/div[1]/div[5]/input')
        except Exception as ex:
            corpo = """<h1>Erro Evope - Erro ao selecionar o cliente</h1><p>Cliente {} e Grupo {}<br>Favor verificar o config e executar o RPA navamente.</p><p>Caso o erro persista entrar em contato com a Evope.</p><p>Muito obrigado!<br>Equipe Evope</p>""".format(clientes[x+1][0],clientes[x+1][1])
            enviar_email(corpo, 'RPA - Erro Evope - Erro ao selecionar o cliente')
            break
        #nome grupo
        
        sleep(1)
        if(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '3'):
            nmGrupo = ''
            agrupados = clientes[numGrupos.groups[nomesClientes[y]][x]][1].split(',')
            for numAgrupados in range(len(agrupados)):
                nmGrupo = nmGrupo + agrupados[numAgrupados].split('>')[-1] + ' e '
            nmGrupo = nmGrupo[0:-2]
        else:
            nmGrupo = clientes[numGrupos.groups[nomesClientes[y]][x]][1]
            if(nmGrupo.count('>')>=1): nmGrupo = nmGrupo.split('>')[-1]
        if(nmGrupo == 'tudo' or nmGrupo == 'geral'): nmGrupo = nomesClientes[y]
        #imagem do mapa
        sleep(2)
        pyauto.press('space')
        sleep(1)
        if(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[3]/div[1]/h3[2]/label[1]').text.split(':')[0] != '0'):
            # percentual hora mapa
            percentualHoraMapa = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[3]/div[1]/h3[2]/label[3]').text
            pyauto.press('pgup')
            pyauto.press('pageup')
            sleep(10)
            #sleep(15)
            pyauto.moveTo(telaX/2, telaY/2+200)
            pyauto.PAUSE = 1.5
            pyauto.click(button='right')
            pyauto.press('down')
            pyauto.press('enter')
            pyperclip.copy(str(nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-Mapa'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])))
            pyauto.hotkey('Ctrl','v')
            pyauto.press('F4')
            pyauto.hotkey('Ctrl','a')
            pyauto.write('%'+'temp%', interval=0.2)
            pyauto.press('enter')
            sleep(2)
            pyauto.hotkey('Alt','l')
            pyauto.PAUSE= 0.2
            pyauto.hotkey('ctrl','f')
            pyauto.write('Lista Fun', interval=0.02)
            pyauto.press('esc')
            pyauto.scroll(50)
            sleep(0.5)
            imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[3]/div[1]/div[2]/div[2]/div').screenshot_as_png)
            im = Image.open(imageStream)
            im.save(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-Processsos'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png')
            pyauto.hotkey('ctrl','f')
            pyauto.write('Mapa de process', interval=0.02)
            pyauto.press('esc')
            pyauto.PAUSE= 1
            #relatório gerencial processos
            driver.switch_to.window(driver.window_handles[1])
            sleep(1)
            for a in range(10):
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(Keys.BACKSPACE)    
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(Keys.DELETE)
            for a in range(10):
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(Keys.BACKSPACE)    
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(Keys.DELETE)            
            sleep(1)
            driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[1]').send_keys(dtInicio)
            sleep(1)
            driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').send_keys(dtFim)
            #selecionar o grupo
            try:
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[2]/span/div/button').click()
                element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[2]/span/div/ul')
                all_li = element.find_elements(By.TAG_NAME, "li")
                grupoNaoSelecionado = True
                for li in all_li:
                    if(li.text=='Select all'):
                            li.click()
                            sleep(1.5)
                            li.click()
                            break
            except Exception as ex:
                pass
            try:
                if(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '2'):
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/h3[1]').text
                    sleep(0.5)
                    try:
                        #duas vezes pq da primeira vez ele fecha, não sei o pq...
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[3]/span/div/button').click()
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/h3[1]').click()
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[3]/span/div/button').click()
                        sleep(0.5)
                        element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[3]/span/div/ul')
                    except Exception as ex:
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[2]/span/div/button').click()
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/h3[1]').click()
                        driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[2]/span/div/button').click()
                        sleep(0.5)
                        element = driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[2]/span/div/ul')
                    sleep(0.5)
                    all_li = element.find_elements(By.TAG_NAME, "li")
                    for li in all_li:
                        if(li.text.upper().count(clientes[numGrupos.groups[nomesClientes[y]][x]][1].upper())==1):
                                li.click()
                                grupoNaoSelecionado = False
                                break
                elif(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '3'):                 
                    agrupados = clientes[numGrupos.groups[nomesClientes[y]][x]][1].split(',')
                    for li in all_li:                    
                        for numAgrupados in range(len(agrupados)):
                            if(li.text.upper().count('MASTER > ' + agrupados[numAgrupados].upper())==1):
                                li.click()
                                grupoNaoSelecionado = False
                                break
                else:
                    if(clientes[numGrupos.groups[nomesClientes[y]][x]][1] == 'tudo'):
                        for li in all_li:
                            if(li.text.upper().count('MASTER')==1):
                                li.click()
                                grupoNaoSelecionado = False
                                break
                    elif(clientes[numGrupos.groups[nomesClientes[y]][x]][1] == 'geral'):
                        for li in all_li:
                            if(li.text.upper().count('SELECT ALL')==1):
                                li.click()
                                grupoNaoSelecionado = False
                                break
                    else:
                        for li in all_li:
                            if(li.text.upper().count('MASTER > ' + clientes[numGrupos.groups[nomesClientes[y]][x]][1].upper())==1):
                                li.click()
                                grupoNaoSelecionado = False
                                break
            except Exception as ex:
                grupoNaoSelecionado = True
            pyauto.PAUSE = 1
            if(grupoNaoSelecionado):
                qtdeColaboradores = 0
            else:
                try:
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[3]/span/div/button').click()
                except Exception as ex:
                    driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]').click()
                sleep(2)
                paginaCarregada('/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]')
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/a[1]').click()
                paginaCarregada('/html/body/form/div[5]/div/div[1]/div/div[1]/input[2]')
                pyauto.PAUSE= 1
                pyauto.press('esc')
                driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]').click()
                pyauto.press('esc')
                gerencialHoras = driver.find_element(By.XPATH,'/html/body/form/div[5]/div/div[2]/div[1]/div[1]/div[1]/table/tbody/tr[2]/td[1]/label').text
                gerencialHoras = gerencialHoras.split(':')[0]
                if(len(gerencialHoras)>=4):
                    gerencialHoras = str(int(int(gerencialHoras)/1000)*1000)
                elif(len(gerencialHoras)>=3):
                    gerencialHoras = str(int(int(gerencialHoras)/100)*100)
                else:
                    gerencialHoras = str(int(int(gerencialHoras)/10)*10)
                gerencialDiario = driver.find_element(By.XPATH,'/html/body/form/div[5]/div/div[2]/div[1]/div[1]/div[2]/table/tbody/tr[2]/td/label').text
                # gerencialDiario = gerencialDiario.split(':')[0]
                if(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '2'):
                    qtdeColaboradores = 1
                else:    
                    qtdeColaboradores =  int(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[1]/div[1]/table/tbody/tr[2]/td[3]/label').text)
                #tirar os prints
                pyauto.press('pgup')
                pyauto.press('pageup')
                sleep(0.5)
                #teste zica merda
                sleep(2)
                imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[1]').screenshot_as_png)
                im = Image.open(imageStream)
                im.save(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G01'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png')
                pyauto.PAUSE= 0.2
                pyauto.hotkey('ctrl','f')
                pyauto.write('Colaboradores com menos', interval=0.02)
                pyauto.press('esc')
                pyauto.scroll(50)
                pyauto.scroll(50)
                sleep(0.5)
                imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[4]/div/canvas').screenshot_as_png)
                im = Image.open(imageStream)
                im.save(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G02'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png')
                pyauto.hotkey('ctrl','f')
                pyauto.write('Aplicativos mais utilizados', interval=0.02)
                pyauto.press('esc')
                pyauto.scroll(50)
                pyauto.scroll(50)
                sleep(0.5)
                try:
                    imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[7]/div[1]/div[2]').screenshot_as_png)
                except Exception as ex:
                    imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[7]/div[2]/div[2]').screenshot_as_png)
                im = Image.open(imageStream)
                im.save(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G03'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png')
                pyauto.hotkey('ctrl','f')
                pyauto.write('Categorias', interval=0.02)
                pyauto.press('esc')
                sleep(0.5)
                pyauto.press('esc')
                pyauto.press('esc')
                pyauto.scroll(50)
                pyauto.scroll(50)
                sleep(0.5)
                imageStream = io.BytesIO(driver.find_element(By.XPATH, '/html/body/form/div[5]/div/div[2]/div[1]/div[8]/div[1]').screenshot_as_png)
                im = Image.open(imageStream)
                im.save(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G04'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png')
                pyauto.hotkey('ctrl','f')
                pyauto.write('Relatório gerencial', interval=0.02)
                pyauto.press('esc')
                pyauto.scroll(50)
                pyauto.scroll(50)
                pyauto.PAUSE= 1
                
        else:
            qtdeColaboradores = 0
        
        try:
            #criar pastas com o nome dos clientes
            os.mkdir('C:\\apresentacoes\\' + nomesClientes[y].replace(' ',''))
        except Exception as ex:
            pass
        # qtdeColaboradoresGrupo += qtdeColaboradores
        driver.switch_to.window(driver.window_handles[0])
        sleep(2)
        #carregar as demais informações do hidden
        hdnTemp = driver.execute_script('return arguments[0].value;', driver.find_element(By.XPATH, '/html/body/form/div[5]/div/input[1]'))
        hdnTemp = hdnTemp.replace('sep=,\n','')
        hdnTemp = hdnTemp.replace('"','')
        rows = hdnTemp.split('\n')
        tempHdn = []
        for row in rows: 
            tempHdn.append(row.split(","))
        if(len(tempHdn)>6):
            dfFuncoes = pd.DataFrame(tempHdn, columns = ['Aplicativo','Funcionalidade', 'Horas', 'Teclado', 'Mouse', 'Interacoes', 'Colaboradores', 'Copiar', 'Colar', 'Processo'])
            dfFuncoes = dfFuncoes.drop(dfFuncoes.index[0])
            # salvar CSV na pasta
            dfFuncoes.to_excel('C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'/'+ nmGrupo.replace(' ','').replace('/','') + '-Funcoes.xlsx')
            dfFuncoes['Teclado'] = pd.to_numeric(dfFuncoes['Teclado'])
            dfFuncoes['Mouse'] = pd.to_numeric(dfFuncoes['Mouse'])
            dfFuncoes['Interacoes'] = pd.to_numeric(dfFuncoes['Interacoes'])
            dfFuncoes['Copiar'] = pd.to_numeric(dfFuncoes['Copiar'])
            dfFuncoes['Colar'] = pd.to_numeric(dfFuncoes['Colar'])
            dfFuncoes['Processo'] = pd.to_numeric(dfFuncoes['Processo'])
            dfFuncoes['hora'] = 1
            totalHora = 0
            totalMinuto = 0
            for i in range(len(dfFuncoes)):
                totalHora += int(dfFuncoes['Horas'].iloc[i].split(':')[0])
                totalMinuto += int(dfFuncoes['Horas'].iloc[i].split(':')[1])
                dfFuncoes['hora'].iloc[i] = int(dfFuncoes['Horas'].iloc[i].split(':')[0])
            totalHora += int(totalMinuto/60)
            principalHora = 0
            totalMinuto = 0
            #pegar as 5 tarefas com mais tempo
            dfFuncoesPrincipais = dfFuncoes.sort_values(by='hora', ascending=False).head(5)
            for i in range(len(dfFuncoesPrincipais)):
                principalHora += int(dfFuncoesPrincipais['Horas'].iloc[i].split(':')[0])
                totalMinuto += int(dfFuncoesPrincipais['Horas'].iloc[i].split(':')[1])
            principalHora += int(totalMinuto/60)
            if(principalHora >0):
                perTempo = principalHora * 100 / totalHora
            else:
                perTempo = 0
            #pegar as 5 tarefas com mais esforço = teclado
            dfFuncoesPrincipais = dfFuncoes.sort_values(by='Teclado', ascending=False).head(5)
            perEsforco = dfFuncoesPrincipais['Teclado'].sum() * 100 / dfFuncoes['Teclado'].sum()
            #importar json
            hdnTemp = driver.execute_script('return arguments[0].value;', driver.find_element(By.XPATH, '/html/body/form/div[5]/div/input[2]'))
            hdnTemp = hdnTemp.replace('sep=,\n','')
            hdnTemp = hdnTemp.replace('"','')
            rows = hdnTemp.split('\n')
            tempHdn = []
            for row in rows: 
                tempHdn.append(row.split(","))
            dfCaminhos = pd.DataFrame(tempHdn, columns = ['Origem','Horas','Destino','Horas2','Quantidade','Copiar Origem','Colar Destino'])
            dfCaminhos = dfCaminhos.drop(dfCaminhos.index[0:1])
            # salvar CSV na pasta
            dfCaminhos.to_excel('C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'/'+ nmGrupo.replace(' ','').replace('/','') + '-Caminhos.xlsx')
            dfCaminhos['Quantidade'] = pd.to_numeric(dfCaminhos['Quantidade'])
            dfCaminhos['Copiar Origem'] = pd.to_numeric(dfCaminhos['Copiar Origem'])
            dfCaminhos['Colar Destino'] = pd.to_numeric(dfCaminhos['Colar Destino'])
            dfCaminhos = dfCaminhos.sort_values(by='Quantidade', ascending=False)
            analisando = True
            a0 = 0
            a1 = dfCaminhos['Origem'].iloc[a0]
            a2 = dfCaminhos['Destino'].iloc[a0]
            a0 += 1
            while analisando:
                if(dfCaminhos['Origem'].iloc[a0] != a2 or dfCaminhos['Destino'].iloc[a0] != a1):
                    a3 = dfCaminhos['Origem'].iloc[a0]
                    a4 = dfCaminhos['Destino'].iloc[a0]
                    analisando = False
                else:
                    a0 += 1
            analisando = True
            a0 += 1
            while analisando:
                if(dfCaminhos['Origem'].iloc[a0] != a2 or dfCaminhos['Destino'].iloc[a0] != a1):
                    if(dfCaminhos['Origem'].iloc[a0] != a4 or dfCaminhos['Destino'].iloc[a0] != a3):
                        a5 = dfCaminhos['Origem'].iloc[a0]
                        a6 = dfCaminhos['Destino'].iloc[a0]
                        analisando = False
                    else:
                        a0 += 1
                else:
                    a0 += 1
            # Preenchimento PPTX e salvar
            #slide 1
            slide = prs.slides[1+6*x]
            trocarTextoMantendoFormatacaoGerencial(slide.shapes[10].text_frame.paragraphs[0], nmGrupo)
            txtTemp = slide.shapes[11].text_frame.paragraphs[0].text
            txtTemp = txtTemp.replace('[diario]',gerencialDiario)
            trocarTextoMantendoFormatacao(slide.shapes[11].text_frame.paragraphs[0], txtTemp)
            txtTemp = slide.shapes[11].text_frame.paragraphs[1].text
            txtTemp = txtTemp.replace('[horas]',gerencialHoras)
            trocarTextoMantendoFormatacao(slide.shapes[11].text_frame.paragraphs[1], txtTemp)
            #imagem 1
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G01'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[4]
            x1, y1, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height
            if(clientes[numGrupos.groups[nomesClientes[y]][x]][2] == '2'):
                new_picture = slide.shapes.add_picture(mapaNovo, x1, (y1+(old_picture.height/4)), old_picture.width, old_picture.height/2)
            else:
                new_picture = slide.shapes.add_picture(mapaNovo, x1, y1, old_picture.width, old_picture.height)
            old_picture.crop_right = 1
            #imagem 2
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G02'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[6]
            x1, y1, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height
            new_picture = slide.shapes.add_picture(mapaNovo, x1, y1)
            new_picture.height = int(new_picture.height / (new_picture.width / 5861835))
            new_picture.width = 5861835
            if(new_picture.height > 2317171): 
                new_picture.height = 2317171
                diferencaTop = 0
            else:
                new_picture.top = new_picture.top + int((2317171 - new_picture.height)/4/2)
                diferencaTop = int((2317171 - new_picture.height)/4)
            old_picture.crop_right = 1
            #imagem 4
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G04'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[7]
            x1, y1, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height
            if(diferencaTop == 0):
                new_picture = slide.shapes.add_picture(mapaNovo, x1, y1, old_picture.width, old_picture.height)
            else:
                new_picture = slide.shapes.add_picture(mapaNovo, x1, y1)
                new_picture.height = old_picture.height + (diferencaTop*3)
                new_picture.width = old_picture.width
                new_picture.width = int(old_picture.width * (new_picture.height / old_picture.height))
                new_picture.top = new_picture.top - int(diferencaTop*3)
                new_picture.left = new_picture.left - int((new_picture.width - old_picture.width)/2)
                diferencaLeft = int((new_picture.width - old_picture.width)/2)
            old_picture.crop_right = 1
            #imagem 3
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-G03'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[5]
            x1, y1 = old_picture.left, old_picture.top
            new_picture = slide.shapes.add_picture(mapaNovo, x1, y1)
            new_picture.height = int(new_picture.height / (new_picture.width / 2477658))
            new_picture.width = 2477658
            if(new_picture.height > 2610504): 
                new_picture.height = 2610504
            if(diferencaTop != 0):
                new_picture.top = new_picture.top - int(diferencaTop*3/2)
                new_picture.left = new_picture.left - diferencaLeft
            old_picture.crop_right = 1
            #slide 2
            slide = prs.slides[2+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[4].text_frame.paragraphs[0], nmGrupo)
            #trocar imagem de processo
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-Mapa'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[6]
            x1, y1, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height            
            new_picture = slide.shapes.add_picture(mapaNovo, x1, y1, old_picture.width, old_picture.height)
            new_picture.crop_right = 0.08
            new_picture.crop_left = 0.12
            old_picture.crop_right = 1
            #slide 3
            slide = prs.slides[3+6*x]
            #imagem do processo
            mapaNovo = pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-Processsos'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png'
            old_picture = slide.shapes[26]
            x1, y1, cx, cy = old_picture.left, old_picture.top, old_picture.width, old_picture.height
            #new_picture = slide.shapes.add_picture(mapaNovo, x1, y1, old_picture.width, old_picture.height)
            new_picture = slide.shapes.add_picture(mapaNovo, x1, y1)
            new_picture.height = int(new_picture.height / (new_picture.width / 2743200))
            new_picture.width = 2743200
            if(new_picture.height > 2930000): new_picture.height = 2930000
            old_picture.crop_right = 1
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[19].text_frame.paragraphs[0], nmGrupo)
            trocarTextoMantendoFormatacao(slide.shapes[5].text_frame.paragraphs[0], str(qtdeColaboradores))
            trocarTextoMantendoFormatacao(slide.shapes[9].text_frame.paragraphs[0], str(totalHora)+'h')
            trocarTextoMantendoFormatacao(slide.shapes[4].text_frame.paragraphs[0], str(int(dfFuncoes['Interacoes'].sum()/1000))+'K')
            trocarTextoMantendoFormatacao(slide.shapes[20].text_frame.paragraphs[0], str(percentualHoraMapa)+'%'+' Tempo')
            #preencher o % de tempo e esfoço e preencher a tabela com as principais funções
            try:
                # trocarTextoMantendoFormatacao(slide.shapes[24].text_frame.paragraphs[0], str(str(int(perTempo))+'% '+ 'Tempo\n'+str(int(perEsforco))+'% '+'Esforço'))
                trocarTextoMantendoFormatacao(slide.shapes[24].text_frame.paragraphs[0], str(str(int(perEsforco)) + '% '+'Esforço'))
                #tabela
                for a in range(5):
                    trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[a+1].cells[0].text_frame.paragraphs[0], '['+dfFuncoesPrincipais['Aplicativo'].iloc[a]+'] '+dfFuncoesPrincipais['Funcionalidade'].iloc[a])
                    trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[a+1].cells[1].text_frame.paragraphs[0], dfFuncoesPrincipais['Horas'].iloc[a])
                    trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[a+1].cells[2].text_frame.paragraphs[0], str(round(dfFuncoesPrincipais['Teclado'].iloc[a]*100/dfFuncoesPrincipais['Teclado'].sum(),2))+' %')
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[24].text_frame.paragraphs[0], '0% '+'Esforço')
            #slide 4
            slide = prs.slides[4+6*x]
            trocarTextoMantendoFormatacao(slide.shapes[29].text_frame.paragraphs[0], a1.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[14].text_frame.paragraphs[0], a2.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[11].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a1) & (dfCaminhos['Destino']==a2)]['Quantidade'].to_numpy()[0]))
            try:
                tempTeste = a1.split(']-')[1]
                textoFuncao = 'Funcionalidade'
            except Exception as ex:
                tempTeste = a1.split(']-')[0].replace('[','').replace(']','')
                textoFuncao = 'Aplicativo'
            trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
            if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
            try:
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[27].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            try:
                trocarTextoMantendoFormatacao(slide.shapes[12].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a2) & (dfCaminhos['Destino']==a1)]['Quantidade'].to_numpy()[0]))
                try:
                    tempTeste = a2.split(']-')[1]
                    textoFuncao = 'Funcionalidade'
                except Exception as ex:
                    tempTeste = a2.split(']-')[0].replace('[','').replace(']','')
                    textoFuncao = 'Aplicativo'
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
                if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[12].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[15].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            trocarTextoMantendoFormatacao(slide.shapes[31].text_frame.paragraphs[0], a3.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[24].text_frame.paragraphs[0], a4.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[21].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a3) & (dfCaminhos['Destino']==a4)]['Quantidade'].to_numpy()[0]))
            try:
                tempTeste = a3.split(']-')[1]
                textoFuncao = 'Funcionalidade'
            except Exception as ex:
                tempTeste = a3.split(']-')[0].replace('[','').replace(']','')
                textoFuncao = 'Aplicativo'
            trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
            if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
            try:
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[26].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            try:
                trocarTextoMantendoFormatacao(slide.shapes[22].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a4) & (dfCaminhos['Destino']==a3)]['Quantidade'].to_numpy()[0]))
                try:
                    tempTeste = a4.split(']-')[1]
                    textoFuncao = 'Funcionalidade'
                except Exception as ex:
                    tempTeste = a4.split(']-')[0].replace('[','').replace(']','')
                    textoFuncao = 'Aplicativo'
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
                if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[22].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[25].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            #slide 5
            slide = prs.slides[5+6*x]
            trocarTextoMantendoFormatacao(slide.shapes[11].text_frame.paragraphs[0], a5.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[15].text_frame.paragraphs[0], a6.split('-')[0].replace('[','').replace(']',''))
            trocarTextoMantendoFormatacao(slide.shapes[8].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a5) & (dfCaminhos['Destino']==a6)]['Quantidade'].to_numpy()[0]))
            try:
                trocarTextoMantendoFormatacao(slide.shapes[9].text_frame.paragraphs[0], str(dfCaminhos[(dfCaminhos['Origem']==a6) & (dfCaminhos['Destino']==a5)]['Quantidade'].to_numpy()[0]))
                try:
                    tempTeste = a6.split(']-')[1]
                    textoFuncao = 'Funcionalidade'
                except Exception as ex:
                    tempTeste = a6.split(']-')[0].replace('[','').replace(']','')
                    textoFuncao = 'Aplicativo'
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
                if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[9].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[13].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            try:
                tempTeste = a5.split(']-')[1]
                textoFuncao = 'Funcionalidade'
            except Exception as ex:
                tempTeste = a5.split(']-')[0].replace('[','').replace(']','')
                textoFuncao = 'Aplicativo'
            trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[0].text_frame.paragraphs[0], tempTeste)
            if(len(tempTeste) >55): tempTeste = tempTeste[:55] + '...'
            try:
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[1].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Horas'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[2].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Interacoes'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[3].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Copiar'].to_numpy()[0]))
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[4].text_frame.paragraphs[0], str(dfFuncoes[(dfFuncoes[textoFuncao]==tempTeste)]['Colar'].to_numpy()[0]))
            except Exception as ex:
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[1].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[2].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[3].text_frame.paragraphs[0], '0')
                trocarTextoMantendoFormatacao(slide.shapes[12].table.rows[1].cells[4].text_frame.paragraphs[0], '0')
            #tabela ordenada por copiar e tabela tabela ordenada por colar
            dfFuncoesTemp = dfFuncoes.sort_values(by='Colar', ascending=False).head(5)
            for i in range(5):
                if(str(dfFuncoesTemp['Funcionalidade'].iloc[i]) == ''):
                    trocarTextoMantendoFormatacao(slide.shapes[17].table.rows[i+1].cells[0].text_frame.paragraphs[0], '['+str(dfFuncoesTemp['Aplicativo'].iloc[i])+']')    
                else:
                    trocarTextoMantendoFormatacao(slide.shapes[17].table.rows[i+1].cells[0].text_frame.paragraphs[0], '['+str(dfFuncoesTemp['Aplicativo'].iloc[i])+']-'+str(dfFuncoesTemp['Funcionalidade'].iloc[i]))
                trocarTextoMantendoFormatacao(slide.shapes[17].table.rows[i+1].cells[1].text_frame.paragraphs[0], str(dfFuncoesTemp['Colar'].iloc[i]))
            dfFuncoesTemp = dfFuncoes.sort_values(by='Copiar', ascending=False).head(5)
            for i in range(5):
                if(str(dfFuncoesTemp['Funcionalidade'].iloc[i]) == ''):
                    trocarTextoMantendoFormatacao(slide.shapes[16].table.rows[i+1].cells[0].text_frame.paragraphs[0], '['+str(dfFuncoesTemp['Aplicativo'].iloc[i])+']')
                else:
                    trocarTextoMantendoFormatacao(slide.shapes[16].table.rows[i+1].cells[0].text_frame.paragraphs[0], '['+str(dfFuncoesTemp['Aplicativo'].iloc[i])+']-'+str(dfFuncoesTemp['Funcionalidade'].iloc[i]))
                trocarTextoMantendoFormatacao(slide.shapes[16].table.rows[i+1].cells[1].text_frame.paragraphs[0], str(dfFuncoesTemp['Copiar'].iloc[i]))
            #slide 6
            slide = prs.slides[6+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[6].text_frame.paragraphs[0], nmGrupo)
        else:
            slide = prs.slides[1+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[10].text_frame.paragraphs[0], nmGrupo)
            slide = prs.slides[2+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[4].text_frame.paragraphs[0], nmGrupo)
            slide = prs.slides[3+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[19].text_frame.paragraphs[0], nmGrupo)    
            slide = prs.slides[6+6*x]
            trocarTextoMantendoFormatacaoNmGrupo(slide.shapes[6].text_frame.paragraphs[0], nmGrupo)
        
        try:
            # #criar pastas com o nome dos clientes
            # os.mkdir('C:\\apresentacoes\\' + nomesClientes[y].replace(' ',''))
            # # salvar os CSVs
            # dfFuncoes.to_excel('C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'/'+ nmGrupo.replace(' ','').replace('/','') + '-Funcoes.xlsx')
            # dfCaminhos.to_excel('C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'/'+ nmGrupo.replace(' ','').replace('/','') + '-Caminhos.xlsx')
            #copiar imagem do mapa
            os.rename(pastaTemp+'\\'+nomesClientes[y].replace(' ','')+'-'+nmGrupo.replace(' ','').replace('/','')+'-Mapa'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2])+'.png', 'C:/apresentacoes/' + nomesClientes[y].replace(' ','') +'/'+ nmGrupo.replace(' ','').replace('/','') + '-Mapa.png')
        except Exception as ex:
            pass
            
    #deixar para o final do grupo editar a página 1 para somar o total de colaboradores
    slide = prs.slides[0]
    trocarTextoMantendoFormatacao(slide.shapes[1].text_frame.paragraphs[0], slide.shapes[1].text_frame.paragraphs[0].text.replace('nmCliente', nomesClientes[y]))
    tempTexto = slide.shapes[1].text_frame.paragraphs[1].text.replace('dtInicio', dtInicio.split('/')[0]+' de '+mesesPortugues[int(dtInicio.split('/')[1])])
    tempTexto = tempTexto.replace('dtFim', dtFim.split('/')[0]+' de '+mesesPortugues[int(dtFim.split('/')[1])]+' de '+dtFim.split('/')[2])
    trocarTextoMantendoFormatacao(slide.shapes[1].text_frame.paragraphs[1], tempTexto)
    trocarTextoMantendoFormatacao(slide.shapes[1].text_frame.paragraphs[2], slide.shapes[1].text_frame.paragraphs[2].text.replace('qtdeColab', str(qtdeColaboradoresGrupo)))
    trocarTextoMantendoFormatacao(slide.shapes[1].text_frame.paragraphs[3], slide.shapes[1].text_frame.paragraphs[3].text.replace('qtdeGrupos', str(len(numGrupos.groups[nomesClientes[y]]))))
    
    #salvar com o nome do cliente e a data
    prs.save('C:/apresentacoes/' + nomesClientes[y].replace(' ','')+'-'+str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2]) + '.pptx')
    #remover os slides que o grupo for zero
    for i in range(len(prs.slides)-5, 1, -6):
        slide = prs.slides[i]
        if(slide.shapes[9].text_frame.paragraphs[0].text == '000h'):
            for j in range(-3,3,1):
                rId = prs.slides._sldIdLst[i-j].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i-j]
                sleep(0.3)
                
    #salvar com o nome do cliente e a data
    prs.save('C:/apresentacoes/' + nomesClientes[y].replace(' ','') + '/' + nomesClientes[y].replace(' ','') + '-' + str(dtInicio.split('/')[0])+str(dtInicio.split('/')[1])+str(dtInicio.split('/')[2]) + '.pptx')
driver.close()
driver.quit()
#obter todos os objetos de uma página
#slide0_shapes = [shape for shape in slide.shapes]